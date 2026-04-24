"""
Веб-сервер ассистента (stdlib http.server).

Поднимает HTTP-сервер на локальном порту и отдаёт:
  GET  /                 → webui/index.html
  GET  /style.css, /app.js, /favicon.ico → статика
  GET  /api/info         → {model, version}
  POST /api/chat         → {message} → {reply}
  GET  /api/events       → SSE поток статус-обновлений (что делает ассистент)

Статус-события извлекаются из перехваченного stdout: строки вида
[tool_name(...)] и [AgentName] ... превращаются в понятные пользователю фразы.

Запуск:
  python main.py --web               # сервер + окно (Edge --app) / браузер
  python main.py --web --no-open     # только сервер
  python main.py --web --port 8765
"""
from __future__ import annotations

import base64
import io
import json
import os
import queue
import re
import socket
import subprocess
import sys
import threading
import time
import uuid
import webbrowser
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from typing import Optional

_ROOT = os.path.dirname(os.path.abspath(__file__))
_WEBUI_DIR = os.path.join(_ROOT, "webui")
_ATTACH_DIR = os.path.join(_ROOT, "attachments")
os.makedirs(_ATTACH_DIR, exist_ok=True)

# ── Карты: инструмент/агент → человекочитаемая фраза ─────────────────────
_TOOL_STATUS = {
    "tavily_search": "Ищу в интернете…",
    "tavily_extract": "Читаю страницу…",
    "open_app": "Запускаю приложение…",
    "open_url": "Открываю ссылку…",
    "open_bookmark": "Открываю закладку…",
    "search_bookmarks": "Ищу закладку…",
    "list_bookmarks_browsers": "Смотрю браузеры…",
    "read_file": "Читаю файл…",
    "write_file": "Записываю файл…",
    "list_directory": "Смотрю папку…",
    "delete_file": "Удаляю файл…",
    "get_weather": "Смотрю погоду…",
    "control_volume": "Настраиваю громкость…",
    "control_media": "Управляю медиа…",
    "browser_navigate": "Перехожу на сайт…",
    "browser_click": "Кликаю на странице…",
    "browser_input_text": "Ввожу текст на сайте…",
    "browser_get_state": "Смотрю страницу…",
    "browser_scroll": "Прокручиваю страницу…",
    "ui_list_windows": "Смотрю открытые окна…",
    "ui_click": "Кликаю по интерфейсу…",
    "ui_click_element": "Нажимаю кнопку…",
    "ui_send_keys": "Нажимаю клавиши…",
    "ui_focus_window": "Переключаюсь на окно…",
    "ui_wait_for_window": "Жду появления окна…",
    "ui_screenshot": "Делаю снимок экрана…",
    "task_done": "Завершаю…",
    # Office / COM
    "office_docs_search": "Ищу в документах…",
    "com_run_python": "Работаю с Office…",
    "office_run_python": "Работаю с Office…",
    "office_close_dialogs": "Закрываю диалоги Office…",
    "office_user_folder": "Смотрю папку пользователя…",
    "office_available_apps": "Смотрю доступные Office-приложения…",
    "office_running_apps": "Смотрю запущенные Office-приложения…",
    "office_is_available": "Проверяю Office-приложение…",
    "office_launch": "Запускаю Office-приложение…",
    "office_quit": "Закрываю Office-приложение…",
    "office_visible": "Меняю видимость окна Office…",
    # Excel
    "excel_get_sheets": "Смотрю листы Excel…",
    "excel_read_sheet": "Читаю лист Excel…",
    "excel_write_cell": "Записываю ячейку в Excel…",
    "excel_write_range": "Записываю диапазон в Excel…",
    "excel_apply_formula": "Применяю формулу в Excel…",
    "excel_create_workbook": "Создаю книгу Excel…",
    # Word
    "word_create_document": "Создаю документ Word…",
    "word_read_document": "Читаю документ Word…",
    "word_write_text": "Пишу в документ Word…",
    "word_find_replace": "Ищу и заменяю в Word…",
    "word_get_tables": "Смотрю таблицы Word…",
    # PowerPoint
    "ppt_create": "Создаю презентацию…",
    "ppt_add_slide": "Добавляю слайд…",
    "ppt_add_textbox": "Добавляю текст на слайд…",
    "ppt_read_slides": "Читаю слайды…",
    # Outlook
    "outlook_send_mail": "Отправляю письмо…",
    "outlook_list_inbox": "Смотрю входящие…",
}
_AGENT_STATUS = {
    "HostAgent": "Планирую задачу…",
    "BrowserAgent": "Работаю с браузером…",
    "WebAgent": "Ищу в интернете…",
    "SystemAgent": "Работаю с системой…",
    "VisionAgent": "Смотрю на экран…",
    "ChatAgent": "Думаю…",
}
_TOOL_RE = re.compile(r"^\s*\[([a-zA-Z_][\w]*)\((.*)\)\]\s*$")
_AGENT_RE = re.compile(r"^\s*\[([A-Z][A-Za-z]+Agent)\]\s+(.+)$")


def _tool_phrase(name: str) -> str:
    if name in _TOOL_STATUS:
        return _TOOL_STATUS[name]
    if name.startswith("browser_"):
        return "Работаю с браузером…"
    if name.startswith("ui_"):
        return "Работаю с интерфейсом…"
    if name.startswith("excel_"):
        return "Работаю с Excel…"
    if name.startswith("word_"):
        return "Работаю с Word…"
    if name.startswith("ppt_"):
        return "Работаю с PowerPoint…"
    if name.startswith("outlook_"):
        return "Работаю с Outlook…"
    if name.startswith("office_"):
        return "Работаю с Office…"
    return f"Выполняю: {name}…"


# ── Брокер статус-событий для SSE ────────────────────────────────────────
class EventBroker:
    def __init__(self) -> None:
        self._subs: list[queue.Queue] = []
        self._lock = threading.Lock()
        self._last: Optional[str] = None

    def subscribe(self) -> queue.Queue:
        q: queue.Queue = queue.Queue(maxsize=128)
        with self._lock:
            self._subs.append(q)
        return q

    def unsubscribe(self, q: queue.Queue) -> None:
        with self._lock:
            if q in self._subs:
                self._subs.remove(q)

    def publish(self, text: str) -> None:
        if text == self._last:
            return
        self._last = text
        self.emit("status", {"text": text})

    def emit(self, event: str, data: dict) -> None:
        payload = json.dumps({"event": event, "data": data}, ensure_ascii=False)
        with self._lock:
            subs = list(self._subs)
        for q in subs:
            try:
                q.put_nowait(payload)
            except queue.Full:
                pass


BROKER = EventBroker()


# ── Перехват stdout → публикация статусов ────────────────────────────────
class _StdoutTap(io.TextIOBase):
    def __init__(self, original) -> None:
        self._orig = original
        self._buf = ""

    def write(self, s: str) -> int:
        try:
            self._orig.write(s)
            self._orig.flush()
        except Exception:
            pass
        self._buf += s
        while "\n" in self._buf:
            line, self._buf = self._buf.split("\n", 1)
            self._dispatch(line)
        return len(s)

    def flush(self) -> None:
        try:
            self._orig.flush()
        except Exception:
            pass

    @staticmethod
    def _dispatch(line: str) -> None:
        if not line.strip():
            return
        m = _TOOL_RE.match(line)
        if m:
            BROKER.publish(_tool_phrase(m.group(1)))
            return
        m = _AGENT_RE.match(line)
        if m:
            BROKER.publish(_AGENT_STATUS.get(m.group(1), f"{m.group(1)} работает…"))


def install_stdout_tap() -> None:
    if not isinstance(sys.stdout, _StdoutTap):
        sys.stdout = _StdoutTap(sys.stdout)


# ── HostAgent (ленивая инициализация) ────────────────────────────────────
_HOST = None
_HOST_LOCK = threading.Lock()
_HOST_READY = threading.Event()


def _ensure_host_async() -> None:
    def _worker():
        global _HOST
        try:
            import main as _m
            _m._ensure_apps_scanned()
            _m._start_ws_bridge()
            with _HOST_LOCK:
                _HOST = _m._make_host_agent()
            BROKER.publish("")
            _HOST_READY.set()
        except Exception as e:
            BROKER.publish(f"Ошибка инициализации: {e}")
            print(f"[web] init error: {e}", file=sys.__stderr__)

    threading.Thread(target=_worker, daemon=True).start()


def _build_chat_history(conv_id: Optional[int], skip_last_user: bool = True,
                        max_pairs: int = 6, max_chars: int = 4000) -> str:
    """Собирает краткую историю чата для подстановки в контекст LLM.
    Берёт последние max_pairs пар user↔assistant, ограничивает общий объём."""
    if not conv_id:
        return ""
    try:
        import database as _db
        items = _db.msg_list(conv_id)
    except Exception:
        return ""
    if skip_last_user and items and items[-1].get("role") == "user":
        items = items[:-1]
    # Last max_pairs*2 messages
    items = items[-(max_pairs * 2):]
    lines: list = []
    for m in items:
        role = "Пользователь" if m.get("role") == "user" else "Ассистент"
        content = (m.get("content") or "").strip()
        if not content:
            continue
        lines.append(f"{role}: {content}")
    text = "\n".join(lines)
    if len(text) > max_chars:
        text = "…\n" + text[-max_chars:]
    return text


def _dispatch(message: str, conv_id: Optional[int] = None) -> dict:
    if not _HOST_READY.wait(timeout=120):
        return {"voice": "[Ошибка] HostAgent не инициализирован.", "screen": {"blocks": []}}
    import main as _m
    from ui_automation import cancel as _cancel
    from ui_automation import sources as _sources
    scope_key = f"conv:{conv_id}" if conv_id else None
    _cancel.set_scope(scope_key)
    _cancel.clear()
    _sources.reset()
    try:
        windows_ctx = _m._get_windows_context()
        rag_ctx = _m._rag_retrieve(message)
        scenario = _m._match_scenario(message)
        chat_ctx = _build_chat_history(conv_id, skip_last_user=True)
        hint = ""
        if scenario:
            hint += (
                f"[Сценарий пользователя — \"{scenario['name']}\"]\n"
                f"{scenario['body'].strip()}\n"
                "(Следуй шагам сценария, адаптируя под текущее состояние системы.)\n"
            )
        if windows_ctx:
            hint += f"[Открытые окна]\n{windows_ctx}\n"
        if rag_ctx:
            hint += f"[Релевантный опыт]\n{rag_ctx}"
        # chat_ctx идёт ОТДЕЛЬНЫМ параметром — его должен видеть именно Planner,
        # а не только worker'ы (иначе планировщик не поймёт «продолжи», «ещё раз» и т.п.).
        result = _HOST.dispatch(message, context_hint=hint, conv_id=conv_id,
                                chat_history=chat_ctx)
        if hasattr(result, "to_dict"):
            out = result.to_dict()
        else:
            out = {"voice": str(result), "screen": {"blocks": []}}

        # Если агенты использовали веб-поиск — добавим блок «Источники».
        urls = _sources.collect()
        if urls:
            blocks = out.setdefault("screen", {"blocks": []}).setdefault("blocks", [])
            # Не дублируем, если форматтер уже вернул links-блок с этими URL.
            existing = set()
            for b in blocks:
                if b.get("type") == "links":
                    existing.update(b.get("links") or [])
            extra = [u for u in urls if u not in existing][:8]
            if extra:
                blocks.append({"type": "links", "title": "Источники", "links": extra})
        return out
    except Exception as e:
        return {"voice": f"[Ошибка] {e}", "screen": {"blocks": []}}
    finally:
        _cancel.clear_scope()


# ── HTTP handler ─────────────────────────────────────────────────────────
_STATIC_FILES = {
    "/": ("index.html", "text/html; charset=utf-8"),
    "/index.html": ("index.html", "text/html; charset=utf-8"),
    "/style.css": ("style.css", "text/css; charset=utf-8"),
    "/app.js": ("app.js", "application/javascript; charset=utf-8"),
    "/vault.js": ("vault.js", "application/javascript; charset=utf-8"),
    "/bridge.js": ("bridge.js", "application/javascript; charset=utf-8"),
    "/icon.svg": ("../src/Icon_Compass.svg", "image/svg+xml"),
    "/favicon.svg": ("../src/Icon_Compass.svg", "image/svg+xml"),
    "/favicon.ico": ("../src/Icon_Compass.ico", "image/x-icon"),
    "/icon.ico": ("../src/Icon_Compass.ico", "image/x-icon"),
}

_MIME_BY_EXT = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
}


_DOC_EXTS = {
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".rtf",
    ".txt", ".md", ".csv", ".tsv", ".json", ".xml",
    ".yaml", ".yml", ".html", ".htm", ".log",
}
_TEXT_EXTS = {".txt", ".md", ".csv", ".tsv", ".json", ".xml",
              ".yaml", ".yml", ".html", ".htm", ".log", ".rtf"}
_MAX_INLINE_PREVIEW = 8000  # chars per document included in agent hint


def _save_document(data_url_or_b64: str, filename_hint: str) -> Optional[dict]:
    """Сохраняет документ, сохраняя оригинальное имя и расширение."""
    s = data_url_or_b64
    mime = "application/octet-stream"
    if s.startswith("data:"):
        head, _, b64 = s.partition(",")
        m = re.match(r"data:([^;]+);base64", head)
        if m:
            mime = m.group(1)
    else:
        b64 = s
    try:
        raw = base64.b64decode(b64)
    except Exception:
        return None
    base_name = os.path.basename(filename_hint or "file")
    ext = os.path.splitext(base_name)[1].lower() or ""
    if ext not in _DOC_EXTS:
        # неизвестные расширения не принимаем, чтобы не плодить мусор
        return None
    stored = f"{uuid.uuid4().hex}{ext}"
    path = os.path.join(_ATTACH_DIR, stored)
    with open(path, "wb") as f:
        f.write(raw)
    return {
        "url": f"/attachments/{stored}",
        "name": base_name,
        "mime": mime,
        "kind": "doc",
        "path": os.path.abspath(path),
    }


def _extract_doc_text(path: str, full: bool = False) -> str:
    """Best-effort извлечение текста для популярных форматов. Пусто, если не получилось.
    full=True — без обрезки по _MAX_INLINE_PREVIEW (для сохранения в vault)."""
    ext = os.path.splitext(path)[1].lower()
    limit = float("inf") if full else _MAX_INLINE_PREVIEW
    try:
        if ext in _TEXT_EXTS:
            for enc in ("utf-8", "cp1251", "latin-1"):
                try:
                    with open(path, "r", encoding=enc) as f:
                        return f.read() if full else f.read(_MAX_INLINE_PREVIEW + 1)
                except UnicodeDecodeError:
                    continue
            return ""
        if ext == ".pdf":
            try:
                import fitz  # type: ignore  # pymupdf
                parts = []
                total = 0
                with fitz.open(path) as doc:
                    for page in doc:
                        t = page.get_text("text") or ""
                        parts.append(t)
                        total += len(t)
                        if total > limit:
                            break
                text = "\n".join(parts)
                if text.strip():
                    return text
            except Exception:
                pass
            try:
                from pypdf import PdfReader  # type: ignore
            except Exception:
                try:
                    from PyPDF2 import PdfReader  # type: ignore
                except Exception:
                    return ""
            reader = PdfReader(path)
            parts = []
            total = 0
            for page in reader.pages:
                t = page.extract_text() or ""
                parts.append(t)
                total += len(t)
                if total > limit:
                    break
            return "\n".join(parts)
        if ext == ".docx":
            try:
                from docx import Document  # type: ignore
            except Exception:
                return ""
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == ".xlsx":
            try:
                from openpyxl import load_workbook  # type: ignore
            except Exception:
                return ""
            wb = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# Лист: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append("\t".join("" if v is None else str(v) for v in row))
                    if sum(len(l) for l in lines) > limit:
                        break
                if sum(len(l) for l in lines) > limit:
                    break
            return "\n".join(lines)
        if ext == ".pptx":
            try:
                from pptx import Presentation  # type: ignore
            except Exception:
                return ""
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"# Слайд {i}")
                for sh in slide.shapes:
                    if hasattr(sh, "text") and sh.text:
                        parts.append(sh.text)
            return "\n".join(parts)
    except Exception:
        return ""
    return ""


def _save_attachment(data_url_or_b64: str, filename_hint: str = "") -> Optional[dict]:
    """Принимает data:image/...;base64,... или просто base64. Возвращает {url, name}."""
    s = data_url_or_b64
    mime = "image/png"
    if s.startswith("data:"):
        head, _, b64 = s.partition(",")
        m = re.match(r"data:([^;]+);base64", head)
        if m:
            mime = m.group(1)
    else:
        b64 = s
    try:
        raw = base64.b64decode(b64)
    except Exception:
        return None
    ext = {v: k for k, v in _MIME_BY_EXT.items()}.get(mime, ".png")
    name = f"{uuid.uuid4().hex}{ext}"
    path = os.path.join(_ATTACH_DIR, name)
    with open(path, "wb") as f:
        f.write(raw)
    return {"url": f"/attachments/{name}", "name": filename_hint or name, "mime": mime}


class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):  # silence default logs
        return

    # ── helpers ────────────────────────────────────────────────────────
    def _send_json(self, code: int, obj) -> None:
        data = json.dumps(obj, ensure_ascii=False).encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(data)

    def _send_src_asset(self, rel: str) -> None:
        src_root = os.path.abspath(os.path.join(_ROOT, "src"))
        target = os.path.abspath(os.path.join(src_root, rel.replace("/", os.sep)))
        if not target.startswith(src_root + os.sep) and target != src_root:
            self.send_error(403); return
        try:
            with open(target, "rb") as f:
                data = f.read()
        except (FileNotFoundError, IsADirectoryError, OSError):
            self.send_error(404); return
        ext = os.path.splitext(target)[1].lower()
        ctype = {
            ".svg": "image/svg+xml",
            ".png": "image/png",
            ".ico": "image/x-icon",
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".webp": "image/webp", ".gif": "image/gif",
        }.get(ext, "application/octet-stream")
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "public, max-age=3600")
        self.end_headers()
        self.wfile.write(data)

    def _send_static(self, filename: str, content_type: str) -> None:
        path = os.path.normpath(os.path.join(_WEBUI_DIR, filename))
        try:
            with open(path, "rb") as f:
                data = f.read()
        except FileNotFoundError:
            self.send_error(404)
            return
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(data)

    # ── GET ────────────────────────────────────────────────────────────
    def do_GET(self):
        path = self.path.split("?", 1)[0]

        if path in _STATIC_FILES:
            fname, ctype = _STATIC_FILES[path]
            self._send_static(fname, ctype)
            return

        if path.startswith("/src/"):
            self._send_src_asset(path[len("/src/"):])
            return

        if path == "/api/info":
            from ui_automation import llm_config as _llm
            cfg = _llm.get()
            self._send_json(200, {
                "model": cfg["model"],
                "provider": cfg["provider"],
                "base_url": cfg["base_url"],
                "ready": _HOST_READY.is_set(),
            })
            return

        if path == "/api/config":
            from ui_automation import llm_config as _llm
            import database as _db
            cfg = _llm.get()
            providers = [{"id": k, "label": v["label"], "base_url": v["base_url"],
                          "api_key_set": _db.provider_key_has(k)}
                         for k, v in _llm.PROVIDERS.items()]
            safe_cfg = dict(cfg)
            safe_cfg["api_key_set"] = bool(cfg.get("api_key"))
            safe_cfg.pop("api_key", None)
            vp = (cfg.get("vision_provider") or "").strip()
            safe_cfg["vision_api_key_set"] = _db.provider_key_has(vp) if vp else safe_cfg["api_key_set"]
            self._send_json(200, {"config": safe_cfg, "providers": providers})
            return

        if path == "/api/models":
            from ui_automation import llm_config as _llm
            qs = self.path.split("?", 1)[1] if "?" in self.path else ""
            base = None
            api_key = None
            import urllib.parse as _up
            for part in qs.split("&"):
                if part.startswith("base_url="):
                    base = _up.unquote(part[len("base_url="):])
                elif part.startswith("api_key="):
                    api_key = _up.unquote(part[len("api_key="):])
            groups = _llm.list_model_groups(base_url=base, api_key=api_key)
            if groups:
                # groups[i]["models"] — список dict'ов {id, vision}
                flat = [m for g in groups for m in g["models"]]
                self._send_json(200, {"models": flat, "groups": groups})
            else:
                models = _llm.list_models(base_url=base, api_key=api_key)
                self._send_json(200, {"models": models})
            return

        if path == "/api/events":
            self._sse_stream()
            return

        if path == "/api/conversations":
            import database as _db
            self._send_json(200, {"items": _db.conv_list()})
            return

        if path.startswith("/api/conversations/") and path.endswith("/messages"):
            try:
                cid = int(path.split("/")[3])
            except Exception:
                self.send_error(400); return
            import database as _db
            self._send_json(200, {"items": _db.msg_list(cid)})
            return

        if path == "/api/vault/scenarios":
            from ui_automation.rag import vault_manager as _vm
            self._send_json(200, {"items": _vm.list_notes("Scenarios")})
            return

        if path == "/api/vault/documents":
            from ui_automation.rag import vault_manager as _vm
            self._send_json(200, {"items": _vm.list_notes("Attachments")})
            return

        if path.startswith("/api/vault/note/"):
            from ui_automation.rag import vault_manager as _vm
            rel = path[len("/api/vault/note/"):]
            import urllib.parse as _up
            rel = _up.unquote(rel)
            note = _vm.read_note(rel)
            if not note:
                self.send_error(404); return
            self._send_json(200, note)
            return

        if path.startswith("/attachments/"):
            fname = path[len("/attachments/"):]
            safe = os.path.basename(fname)
            fpath = os.path.join(_ATTACH_DIR, safe)
            if not os.path.isfile(fpath):
                self.send_error(404); return
            ext = os.path.splitext(safe)[1].lower()
            ctype = _MIME_BY_EXT.get(ext, "application/octet-stream")
            with open(fpath, "rb") as f:
                data = f.read()
            self.send_response(200)
            self.send_header("Content-Type", ctype)
            self.send_header("Content-Length", str(len(data)))
            self.send_header("Cache-Control", "public, max-age=31536000")
            self.end_headers()
            self.wfile.write(data)
            return

        self.send_error(404)

    def _read_json(self) -> Optional[dict]:
        length = int(self.headers.get("Content-Length", "0") or "0")
        try:
            body = self.rfile.read(length).decode("utf-8") if length else "{}"
            return json.loads(body)
        except Exception:
            return None

    # ── POST ───────────────────────────────────────────────────────────
    def do_POST(self):
        path = self.path.split("?", 1)[0]
        import database as _db

        if path == "/api/chat":
            req = self._read_json()
            if req is None:
                self._send_json(400, {"error": "invalid JSON"}); return
            message = (req.get("message") or "").strip()
            conv_id = req.get("conversation_id")
            images = req.get("images") or []  # list of data-url strings
            documents = req.get("documents") or []  # list of {dataUrl, name, mime}
            if not message and not images and not documents:
                self._send_json(400, {"error": "empty message"}); return

            attachments = []
            doc_attachments = []
            for img in images:
                saved = _save_attachment(img)
                if saved:
                    attachments.append(saved)
            for d in documents:
                if not isinstance(d, dict):
                    continue
                saved = _save_document(d.get("dataUrl") or "", d.get("name") or "file")
                if saved:
                    doc_attachments.append(saved)
                    attachments.append({
                        "url": saved["url"], "name": saved["name"],
                        "mime": saved["mime"], "kind": "doc",
                    })

            # Create conversation if missing
            if not conv_id:
                title = (message or "Изображение")
                conv_id = _db.conv_create(title=title)
                BROKER.emit("conv_created", {"id": conv_id, "title": title})
            elif message:
                # keep title from first user message if empty
                pass

            _db.msg_add(conv_id, "user", message, attachments=attachments or None)
            BROKER.emit("msg_added", {"conversation_id": conv_id, "role": "user",
                                      "content": message, "attachments": attachments})

            # Augment message with attachment hint for agent context
            augmented = message
            img_atts = [a for a in attachments if a.get("kind") != "doc"]
            if img_atts:
                paths = ", ".join(os.path.abspath(os.path.join(_ROOT, a["url"].lstrip("/"))) for a in img_atts)
                augmented = (
                    augmented
                    + f"\n\n[Пользователь приложил к сообщению картинки (не скриншоты экрана): {paths}. "
                    + "Используй vision-агент, чтобы посмотреть именно на эти файлы и ответить по их содержимому. "
                    + "НЕ открывай их в просмотрщике, НЕ делай скриншот экрана.]"
                ).strip()
            if doc_attachments:
                parts = ["\n\n[Прикреплённые документы]"]
                for d in doc_attachments:
                    parts.append(f"— {d['name']} ({d['path']})")
                    preview = _extract_doc_text(d["path"])
                    if preview:
                        truncated = preview[:_MAX_INLINE_PREVIEW]
                        suffix = "\n…[текст обрезан, полный файл доступен по пути выше]" \
                                 if len(preview) > _MAX_INLINE_PREVIEW else ""
                        parts.append(f"--- начало содержимого {d['name']} ---\n"
                                     f"{truncated}{suffix}\n"
                                     f"--- конец содержимого {d['name']} ---")
                augmented = (augmented + "\n" + "\n".join(parts)).strip()

            reply = _dispatch(augmented, conv_id=conv_id)
            from ui_automation import cancel as _cancel
            # Если пользователь нажал «Отмена» — НЕ показываем ответ ассистента вовсе.
            if _cancel.is_cancelled(f"conv:{conv_id}"):
                _cancel.clear(f"conv:{conv_id}")
                self._send_json(200, {"cancelled": True, "conversation_id": conv_id})
                return
            voice = reply.get("voice", "")
            _db.msg_add(conv_id, "assistant", voice, response_json=json.dumps(reply, ensure_ascii=False))
            BROKER.emit("msg_added", {"conversation_id": conv_id, "role": "assistant",
                                      "content": voice, "response": reply})
            BROKER.emit("conv_updated", {"id": conv_id})

            self._send_json(200, {
                "reply": voice, "response": reply,
                "conversation_id": conv_id,
                "attachments": attachments,
            })
            return

        if path == "/api/conversations":
            req = self._read_json() or {}
            title = (req.get("title") or "Новый чат")
            cid = _db.conv_create(title=title)
            BROKER.emit("conv_created", {"id": cid, "title": title})
            self._send_json(200, {"id": cid, "title": title})
            return

        if path == "/api/config":
            from ui_automation import llm_config as _llm
            req = self._read_json() or {}
            cfg = _llm.set_config(
                provider=req.get("provider"),
                model=req.get("model"),
                base_url=req.get("base_url"),
                api_key=req.get("api_key"),
                vision_model=req.get("vision_model"),
                folder=req.get("folder"),
                vision_provider=req.get("vision_provider"),
                vision_base_url=req.get("vision_base_url"),
                vision_api_key=req.get("vision_api_key"),
            )
            safe_cfg = dict(cfg)
            safe_cfg["api_key_set"] = bool(cfg.get("api_key"))
            safe_cfg.pop("api_key", None)
            import database as _db
            vp = (cfg.get("vision_provider") or "").strip()
            safe_cfg["vision_api_key_set"] = _db.provider_key_has(vp) if vp else safe_cfg["api_key_set"]
            BROKER.emit("config_updated", safe_cfg)
            self._send_json(200, {"config": safe_cfg})
            return

        if path == "/api/cancel":
            from ui_automation import cancel as _cancel
            req = self._read_json() or {}
            cid = req.get("conversation_id")
            key = f"conv:{cid}" if cid else None
            _cancel.request_cancel(key)
            BROKER.emit("cancelled", {"conversation_id": cid})
            self._send_json(200, {"ok": True})
            return

        if path == "/api/vault/scenarios":
            from ui_automation.rag import vault_manager as _vm
            req = self._read_json() or {}
            name = (req.get("name") or "").strip()
            body = (req.get("body") or "").strip()
            triggers = req.get("triggers") or []
            tags = req.get("tags") or ["scenario"]
            if not name or not body:
                self._send_json(400, {"error": "name and body are required"}); return
            if isinstance(triggers, str):
                triggers = [t.strip() for t in triggers.split(",") if t.strip()]
            path_saved = _vm.save_scenario(name, triggers, body, tags=tags)
            self._send_json(200, {"ok": True, "path": path_saved})
            return

        if path == "/api/vault/documents":
            from ui_automation.rag import vault_manager as _vm
            req = self._read_json() or {}
            data = req.get("data") or ""
            name_hint = req.get("name") or "file"
            tags_raw = req.get("tags") or []
            if isinstance(tags_raw, str):
                tags = [t.strip() for t in tags_raw.split(",") if t.strip()]
            else:
                tags = list(tags_raw)
            # Сохраняем бинарник в vault/Attachments/
            base_name = os.path.basename(name_hint)
            ext = os.path.splitext(base_name)[1].lower()
            if ext not in _DOC_EXTS:
                self._send_json(400, {"error": "unsupported extension"}); return
            s = data
            if s.startswith("data:"):
                _, _, b64 = s.partition(",")
            else:
                b64 = s
            try:
                raw = base64.b64decode(b64)
            except Exception:
                self._send_json(400, {"error": "bad base64"}); return
            from ui_automation.rag.vault_manager import VAULT_DIR
            import re as _re
            safe = _re.sub(r"[^\w.\- а-яА-ЯёЁ]+", "_", base_name)
            target = os.path.join(VAULT_DIR, "Attachments", safe)
            # не перезаписываем — при коллизии добавляем суффикс
            if os.path.exists(target):
                stem, ex = os.path.splitext(safe)
                target = os.path.join(VAULT_DIR, "Attachments", f"{stem}-{int(time.time())}{ex}")
            with open(target, "wb") as f:
                f.write(raw)
            text = _extract_doc_text(target, full=True)
            md_path = _vm.save_document(target, text or "", tags=tags or ["document"])
            self._send_json(200, {"ok": True, "path": target, "note": md_path,
                                  "preview": (text or "")[:400]})
            return

        if path == "/api/upload":
            req = self._read_json() or {}
            data = req.get("data") or ""
            name = req.get("name") or ""
            saved = _save_attachment(data, name)
            if not saved:
                self._send_json(400, {"error": "bad image"}); return
            self._send_json(200, saved)
            return

        self.send_error(404)

    def do_DELETE(self):
        path = self.path.split("?", 1)[0]
        import database as _db
        if path.startswith("/api/vault/note/"):
            from ui_automation.rag import vault_manager as _vm
            import urllib.parse as _up
            rel = _up.unquote(path[len("/api/vault/note/"):])
            ok = _vm.delete_note(rel)
            self._send_json(200 if ok else 404, {"ok": ok})
            return
        if path.startswith("/api/conversations/"):
            try:
                cid = int(path.split("/")[3])
            except Exception:
                self.send_error(400); return
            _db.conv_delete(cid)
            BROKER.emit("conv_deleted", {"id": cid})
            self._send_json(200, {"ok": True})
            return
        self.send_error(404)

    # ── SSE ────────────────────────────────────────────────────────────
    def _sse_stream(self):
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream")
        self.send_header("Cache-Control", "no-cache")
        self.send_header("Connection", "keep-alive")
        self.send_header("X-Accel-Buffering", "no")
        self.end_headers()

        q = BROKER.subscribe()
        try:
            self.wfile.write(b": connected\n\n")
            self.wfile.flush()
            last_ping = time.time()
            while True:
                try:
                    payload = q.get(timeout=15)
                    try:
                        obj = json.loads(payload)
                        evt = obj.get("event", "status")
                        data_str = json.dumps(obj.get("data", {}), ensure_ascii=False)
                    except Exception:
                        evt, data_str = "status", payload
                    msg = f"event: {evt}\ndata: {data_str}\n\n".encode("utf-8")
                    self.wfile.write(msg)
                    self.wfile.flush()
                except queue.Empty:
                    # keepalive
                    if time.time() - last_ping > 10:
                        self.wfile.write(b": ping\n\n")
                        self.wfile.flush()
                        last_ping = time.time()
        except (BrokenPipeError, ConnectionResetError, ConnectionAbortedError):
            pass
        finally:
            BROKER.unsubscribe(q)


# ── Окно (Edge --app) ────────────────────────────────────────────────────
def _find_edge() -> Optional[str]:
    candidates = [
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


def _open_window(url: str, standalone: bool = True) -> None:
    """Открывает URL в chromeless-окне Edge (--app) или, если нет, в браузере."""
    if standalone:
        edge = _find_edge()
        if edge:
            try:
                subprocess.Popen(
                    [edge, f"--app={url}", "--new-window",
                     "--window-size=900,700"],
                    shell=False,
                )
                return
            except Exception:
                pass
    try:
        webbrowser.open(url, new=1)
    except Exception:
        pass


def _port_free(port: int) -> bool:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        try:
            s.bind(("127.0.0.1", port))
            return True
        except OSError:
            return False


def _pick_port(preferred: int) -> int:
    if _port_free(preferred):
        return preferred
    # Linear scan
    for p in range(preferred + 1, preferred + 20):
        if _port_free(p):
            return p
    return preferred  # will fail and show error


def _start_server(port: int) -> tuple[ThreadingHTTPServer, str]:
    install_stdout_tap()
    _ensure_host_async()
    port = _pick_port(port)
    httpd = ThreadingHTTPServer(("127.0.0.1", port), Handler)
    url = f"http://127.0.0.1:{port}/"
    print(f"[web] Компас запущен: {url}")
    return httpd, url


def run_web(port: int = 8765, open_window: bool = True, standalone: bool = True) -> None:
    httpd, url = _start_server(port)
    if open_window:
        threading.Timer(0.4, lambda: _open_window(url, standalone)).start()
    try:
        httpd.serve_forever()
    except KeyboardInterrupt:
        print("\n[web] Остановлено.")
        httpd.shutdown()


def run_app(port: int = 8765, width: int = 980, height: int = 740) -> None:
    """Запускает UI как десктопное приложение в окне pywebview (Edge WebView2).

    Сервер поднимается в фоновом потоке, а главный поток держит окно webview —
    это требование pywebview (GUI-цикл должен быть в основном потоке).
    """
    try:
        import webview  # type: ignore
    except ImportError:
        print("[web] Нужен пакет pywebview. Установи: pip install pywebview")
        print("[web] Откатываюсь на окно Edge --app.")
        run_web(port=port, open_window=True, standalone=True)
        return

    # Для Windows: задать AppUserModelID и иконку процесса,
    # чтобы панель задач показывала иконку Компаса, а не python.exe.
    try:
        import ctypes
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID("Compass.Assistant.App.1")
    except Exception:
        pass

    httpd, url = _start_server(port)

    # Запускаем WS-мост, чтобы webview-окно подключилось к нему как
    # «расширение» и отвечало на browser_* команды.
    try:
        import browser_extension.ws_server as _ws
        if not _ws.is_running():
            _ws.start_thread()
    except Exception as e:
        print(f"[web] WS bridge start failed: {e}")

    server_thread = threading.Thread(
        target=httpd.serve_forever, daemon=True, name="webui-http"
    )
    server_thread.start()

    app_url = url + "?bridge=1"
    try:
        window = webview.create_window(
            "Компас",
            app_url,
            width=width,
            height=height,
            min_size=(520, 480),
            text_select=True,
        )
        icon_path = os.path.join(_ROOT, "src", "Icon_Compass.ico")

        state = {"exiting": False, "tray": None}

        def on_closing():
            # Пока пользователь явно не выбрал «Выйти» — прячем окно в трей.
            if state["exiting"]:
                return True
            try:
                window.hide()
            except Exception:
                pass
            return False

        try:
            window.events.closing += on_closing
        except Exception:
            pass

        def _show_window():
            try:
                window.show()
                window.restore()
            except Exception:
                pass

        def _open_in_browser():
            try:
                import webbrowser
                webbrowser.open(url)
            except Exception:
                pass

        def _stop_tasks():
            try:
                from ui_automation import cancel as _cancel
                _cancel.request_cancel_all()
            except Exception as e:
                print(f"[tray] cancel failed: {e}")

        def _exit_app():
            state["exiting"] = True
            try:
                if state["tray"] is not None:
                    state["tray"].stop()
            except Exception:
                pass
            try:
                window.destroy()
            except Exception:
                pass

        def _load_tray_image():
            from PIL import Image
            png_path = os.path.join(_ROOT, "src", "png", "orange", "compass-orange-tiny.png")
            for p in (png_path, icon_path):
                if os.path.isfile(p):
                    try:
                        im = Image.open(p)
                        im.load()
                        return im.convert("RGBA").resize((64, 64), Image.LANCZOS)
                    except Exception:
                        continue
            return Image.new("RGBA", (64, 64), (255, 140, 0, 255))

        tray_ready = threading.Event()

        def _start_tray():
            try:
                import pystray
                img = _load_tray_image()
                menu = pystray.Menu(
                    pystray.MenuItem("Открыть приложение", lambda: _show_window(), default=True),
                    pystray.MenuItem("Открыть в браузере", lambda: _open_in_browser()),
                    pystray.MenuItem("Остановить", lambda: _stop_tasks()),
                    pystray.Menu.SEPARATOR,
                    pystray.MenuItem("Выйти", lambda: _exit_app()),
                )
                tray = pystray.Icon("compass", img, "Компас", menu)
                state["tray"] = tray
                tray.run(setup=lambda ic: (ic.__setattr__("visible", True), tray_ready.set()))
            except Exception as e:
                print(f"[tray] disabled ({e})")
                tray_ready.set()

        threading.Thread(target=_start_tray, daemon=True, name="compass-tray").start()
        tray_ready.wait(timeout=2.0)

        start_kwargs = {}
        if os.path.isfile(icon_path):
            start_kwargs["icon"] = icon_path
        try:
            webview.start(**start_kwargs)
        except TypeError:
            # старые версии pywebview без параметра icon
            webview.start()
    except Exception as e:
        print(f"[web] Ошибка webview ({e}), откатываюсь на окно Edge --app.")
        _open_window(url, standalone=True)
        try:
            server_thread.join()
        except KeyboardInterrupt:
            pass
    finally:
        try:
            if 'state' in locals() and state.get("tray") is not None:
                state["tray"].stop()
        except Exception:
            pass
        try:
            httpd.shutdown()
        except Exception:
            pass
        # Принудительно завершаем процесс: у нас живут фоновые потоки
        # (ws-мост, сканеры, HostAgent-воркеры), часть из них — не daemon,
        # и обычный return из run_web() оставляет процесс висеть в трее-невидимкой.
        if 'state' in locals() and state.get("exiting"):
            os._exit(0)


if __name__ == "__main__":
    run_web()
